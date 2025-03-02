from PIL import Image
import os
import time
from pptx import Presentation
from pptx.util import Cm

# Start-Timer
start_time = time.time()

# Ordner, in dem die Plots gespeichert sind
plot_dir = os.path.join("plots")
powerpoint_path = os.path.join("outputs", "HR_BEM.pptx")

# Mindestauflösung festlegen (z. B. 1920x1080px)
MIN_WIDTH = 1000
MIN_HEIGHT = 1000

# PowerPoint-Präsentation erstellen
presentation = Presentation()

# Seitenverhältnis 16:9 in Zentimetern (33.87 cm x 19.05 cm)
presentation.slide_width = Cm(33.87)
presentation.slide_height = Cm(19.05)

# Alle Plot-Bilder durchlaufen
for filename in sorted(os.listdir(plot_dir)):
    if filename.endswith(".png"):  # Nur PNG-Dateien berücksichtigen
        plot_path = os.path.join(plot_dir, filename)

        # Bildauflösung prüfen
        with Image.open(plot_path) as img:
            width, height = img.size  # Breite und Höhe des Bildes in Pixeln

        # Qualifizierte Bilder (mindestens 1920x1080)
        if width >= MIN_WIDTH and height >= MIN_HEIGHT:
            # Neue Folie hinzufügen mit Layout (Titel + Inhalt)
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])

            # Folienüberschrift hinzufügen
            slide.shapes.title.text = f"Analyse: {os.path.splitext(filename)[0]}"  # Dateiname ohne .png

            # Bildgröße anpassen (festes Seitenverhältnis, keine Verzerrung)
            left = Cm(1.27)  # Margen links
            top = Cm(3.81)  # Margen oben
            width_cm = Cm(15.24)
            height_cm = Cm(10.16)
            slide.shapes.add_picture(plot_path, left, top, width=width_cm, height=height_cm)

            # Textfeld neben dem Bild hinzufügen
            text_left = Cm(17.27)
            text_top = Cm(3.81)
            text_width = Cm(11.43)
            text_height = Cm(10.16)
            textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = textbox.text_frame

            # Platzhaltertext für das Textfeld
            text_frame.text = "Hier können Sie die Analyseergebnisse oder zusätzliche Informationen zu diesem Bild einfügen."

            # Schriftart und -größe festlegen
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Cm(0.5)  # Schriftgröße ca. 14 Pt
                paragraph.font.name = "Arial"  # Schriftart
        else:
            print(f"Bild {filename} übersprungen (Auflösung {width}x{height} zu niedrig).")

# PowerPoint-Datei speichern
presentation.save(powerpoint_path)
print(f"PowerPoint-Präsentation wurde gespeichert: {powerpoint_path}")

# Dauerberechnung
end_time = time.time()
print(f"Analyse abgeschlossen in {end_time - start_time:.2f} Sekunden.")